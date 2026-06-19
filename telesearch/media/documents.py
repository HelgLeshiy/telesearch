"""Extract searchable text from attached files and documents.

Telegram messages can carry arbitrary file attachments. This module turns the
*content* of text-bearing files into plain text so they can be embedded and
searched like any other message:

  * PDF                      -> pypdf
  * Word (.docx)             -> python-docx
  * Excel (.xlsx)            -> openpyxl
  * PowerPoint (.pptx)       -> python-pptx
  * HTML                     -> stdlib tag stripping
  * Plain/structured text    -> read directly (txt, md, csv, tsv, json, logs,
                                subtitles, source code, ...)

Truly binary files (zip, exe, images, audio, fonts, ...) yield no text and are
skipped. Optional parsers are imported lazily; if a dependency is missing the
file is skipped with a warning rather than crashing the whole index build.
"""

from __future__ import annotations

import html
import re
from html.parser import HTMLParser
from pathlib import Path
from typing import Optional

# Extensions we read directly as text (covers most "non-binary" attachments).
_TEXT_EXTENSIONS = {
    ".txt", ".text", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json",
    ".jsonl", ".ndjson", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".log", ".srt", ".vtt", ".sub", ".tex", ".xml", ".rtf",
    # source code
    ".py", ".pyi", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".h", ".cpp",
    ".hpp", ".cc", ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".kts",
    ".scala", ".sh", ".bash", ".zsh", ".ps1", ".sql", ".r", ".m", ".lua",
    ".pl", ".dart", ".vue", ".css", ".scss", ".less", ".html", ".htm",
}
_PDF_EXT = {".pdf"}
_DOCX_EXT = {".docx"}
_XLSX_EXT = {".xlsx", ".xlsm"}
_PPTX_EXT = {".pptx"}
_HTML_EXT = {".html", ".htm"}


class _TextExtractingParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in {"script", "style"}:
            self._skip = True

    def handle_endtag(self, tag):
        if tag in {"script", "style"}:
            self._skip = False

    def handle_data(self, data):
        if not self._skip and data.strip():
            self._parts.append(data)

    def text(self) -> str:
        return re.sub(r"\s+\n", "\n", " ".join(self._parts)).strip()


def _strip_html(raw: str) -> str:
    parser = _TextExtractingParser()
    parser.feed(html.unescape(raw))
    return parser.text()


def _read_text_file(path: Path, max_chars: int) -> str:
    """Read a text file, rejecting content that looks binary."""
    raw = path.read_bytes()[: max_chars * 4]
    if b"\x00" in raw[:8192]:
        return ""  # NUL bytes -> binary
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        return ""
    # Guard against mostly-unprintable blobs that happened to decode.
    sample = text[:4096]
    if sample:
        printable = sum(c.isprintable() or c in "\n\r\t " for c in sample)
        if printable / len(sample) < 0.85:
            return ""
    if path.suffix.lower() in _HTML_EXT:
        text = _strip_html(text)
    return text[:max_chars].strip()


def _extract_pdf(path: Path, max_chars: int) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        raise RuntimeError("PDF extraction needs pypdf (pip install pypdf)")
    reader = PdfReader(str(path))
    parts: list[str] = []
    total = 0
    for page in reader.pages:
        chunk = page.extract_text() or ""
        parts.append(chunk)
        total += len(chunk)
        if total >= max_chars:
            break
    return "\n".join(parts)[:max_chars].strip()


def _extract_docx(path: Path, max_chars: int) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        raise RuntimeError("DOCX extraction needs python-docx (pip install python-docx)")
    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)[:max_chars].strip()


def _extract_xlsx(path: Path, max_chars: int) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        raise RuntimeError("XLSX extraction needs openpyxl (pip install openpyxl)")
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    total = 0
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                line = " | ".join(cells)
                parts.append(line)
                total += len(line)
        if total >= max_chars:
            break
    wb.close()
    return "\n".join(parts)[:max_chars].strip()


def _extract_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise RuntimeError("PPTX extraction needs python-pptx (pip install python-pptx)")
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)[:max_chars].strip()


def is_supported(file_name: Optional[str], mime_type: Optional[str]) -> bool:
    """Heuristic: would we be able to extract text from this attachment?"""
    ext = Path(file_name or "").suffix.lower()
    if ext in _TEXT_EXTENSIONS or ext in _PDF_EXT or ext in _DOCX_EXT \
            or ext in _XLSX_EXT or ext in _PPTX_EXT:
        return True
    if mime_type:
        if mime_type.startswith("text/"):
            return True
        if mime_type in {
            "application/pdf",
            "application/json",
            "application/xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }:
            return True
    return False


def extract_document_text(
    path: str | Path,
    mime_type: Optional[str] = None,
    file_name: Optional[str] = None,
    max_chars: int = 400_000,
) -> str:
    """Return extracted text for a file, or "" if it is binary/unsupported."""
    path = Path(path)
    if not path.exists():
        return ""
    ext = (Path(file_name).suffix.lower() if file_name else path.suffix.lower())

    if ext in _PDF_EXT or mime_type == "application/pdf":
        return _extract_pdf(path, max_chars)
    if ext in _DOCX_EXT:
        return _extract_docx(path, max_chars)
    if ext in _XLSX_EXT:
        return _extract_xlsx(path, max_chars)
    if ext in _PPTX_EXT:
        return _extract_pptx(path, max_chars)
    if ext in _TEXT_EXTENSIONS or (mime_type and mime_type.startswith("text/")):
        return _read_text_file(path, max_chars)
    # Last resort: try reading as text; binary content is rejected internally.
    return _read_text_file(path, max_chars)


def split_text(text: str, chunk_chars: int = 1200, overlap: int = 150) -> list[str]:
    """Split long text into overlapping windows, preferring whitespace breaks."""
    text = text.strip()
    if not text:
        return []
    if len(text) <= chunk_chars:
        return [text]

    chunks: list[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_chars, n)
        if end < n:
            # Try to break on the last newline/space in the window.
            window = text[start:end]
            brk = max(window.rfind("\n"), window.rfind(" "))
            if brk > chunk_chars * 0.5:
                end = start + brk
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= n:
            break
        start = max(end - overlap, start + 1)
    return chunks
